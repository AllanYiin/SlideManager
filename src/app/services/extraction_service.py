# -*- coding: utf-8 -*-

from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
from typing import List, Optional

from app.core.logging import get_logger

log = get_logger(__name__)


@dataclass
class SlideText:
    page: int
    title: str
    body: str
    all_text: str


class ExtractionService:
    """使用 python-pptx 抽取文字。"""

    def extract(self, pptx_path: Path) -> List[SlideText]:
        from pptx import Presentation

        prs = Presentation(str(pptx_path))
        out: List[SlideText] = []
        for idx, slide in enumerate(prs.slides, start=1):
            texts: List[str] = []
            title = ""
            body_parts: List[str] = []

            for shape in slide.shapes:
                try:
                    if not hasattr(shape, "has_text_frame") or not shape.has_text_frame:
                        continue
                    t = shape.text_frame.text or ""
                    t = t.strip()
                    if not t:
                        continue
                    texts.append(t)

                    # 嘗試抓 title（placeholder title 或第一段）
                    if not title:
                        # python-pptx placeholder 種類不一定穩
                        ph = getattr(shape, "placeholder_format", None)
                        ptype = getattr(getattr(ph, "type", None), "__str__", lambda: "")()
                        # 保守：只要 shape.name 有 Title/標題字樣
                        if "title" in (shape.name or "").lower() or "標題" in (shape.name or ""):
                            title = t.splitlines()[0].strip()
                        elif len(t) <= 80:
                            # 退化：第一個較短的文字當 title
                            title = t.splitlines()[0].strip()
                    else:
                        body_parts.append(t)
                except Exception:
                    continue

            if not title and texts:
                title = texts[0].splitlines()[0].strip()

            all_text = "\n".join(texts).strip()
            body = "\n".join(body_parts).strip()
            out.append(SlideText(page=idx, title=title, body=body, all_text=all_text))

        return out
