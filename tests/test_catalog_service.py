# -*- coding: utf-8 -*-

import sys
import tempfile
import unittest

import warnings
from pathlib import Path


# 讓 unittest 在任何工作目錄下都能找到 src/app
ROOT = Path(__file__).resolve().parents[1]
SRC = ROOT / "src"
if str(SRC) not in sys.path:
    sys.path.insert(0, str(SRC))

from app.services.catalog_service import CatalogService
from app.services.project_store import ProjectStore



def _build_pptx(path: Path, slides: int = 1) -> bool:
    try:
        from pptx import Presentation
    except Exception as exc:
        warnings.warn(f"無法載入 python-pptx，略過建立測試 PPTX：{exc}")
        return False

    prs = Presentation()
    for _ in range(slides):
        prs.slides.add_slide(prs.slide_layouts[6])
    prs.save(str(path))
    return True


class TestCatalogService(unittest.TestCase):
    def test_whitelist_add_remove_update(self):
        with tempfile.TemporaryDirectory() as td:
            store = ProjectStore(Path(td))
            svc = CatalogService(store)

            dirs = svc.add_whitelist_dir(td)
            self.assertIn(str(Path(td).resolve()), dirs)

            entries = svc.update_whitelist_dir(td, enabled=False, recursive=False)
            self.assertEqual(len(entries), 1)
            self.assertFalse(entries[0]["enabled"])
            self.assertFalse(entries[0]["recursive"])

            dirs = svc.remove_whitelist_dir(td)
            self.assertEqual(dirs, [])


    def test_scan_and_mark_indexed(self):
        with tempfile.TemporaryDirectory() as td:
            if sys.platform != "win32":
                warnings.warn("非 Windows 平台，PPTX 讀取結果可能與預期不同")
            root = Path(td)
            pptx_path = root / "demo.pptx"
            if not _build_pptx(pptx_path, slides=2):
                self.skipTest("無法建立 PPTX，略過測試")


            store = ProjectStore(root)
            svc = CatalogService(store)
            svc.add_whitelist_dir(str(root))

            out = svc.scan()
            files = out["files"]
            self.assertEqual(len(files), 1)
            self.assertEqual(files[0]["filename"], "demo.pptx")
            self.assertEqual(files[0]["slide_count"], 2)

            svc.mark_indexed(files[0]["abs_path"], slides_count=2)
            cat = store.load_catalog()
            self.assertTrue(cat["files"][0]["indexed"])
            self.assertEqual(cat["files"][0]["slides_count"], 2)


    def test_mark_missing_and_clear(self):
        with tempfile.TemporaryDirectory() as td:
            if sys.platform != "win32":
                warnings.warn("非 Windows 平台，PPTX 讀取結果可能與預期不同")
            root = Path(td)
            pptx_path = root / "remove_me.pptx"
            if not _build_pptx(pptx_path, slides=1):
                self.skipTest("無法建立 PPTX，略過測試")


            store = ProjectStore(root)
            svc = CatalogService(store)
            svc.add_whitelist_dir(str(root))
            svc.scan()

            pptx_path.unlink()
            out = svc.scan()
            self.assertTrue(out["files"][0]["missing"])

            removed = svc.clear_missing_files()
            self.assertEqual(removed, 1)
            self.assertEqual(store.load_catalog()["files"], [])


if __name__ == "__main__":
    unittest.main()
